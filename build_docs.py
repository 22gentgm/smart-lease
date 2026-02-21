from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE = os.path.dirname(__file__)
DOCX_PATH = os.path.join(BASE, "True_Fruits_Case_Study.docx")
PPTX_PATH = os.path.join(BASE, "True_Fruits_Presentation.pptx")



def add_simple_table(doc, headers, rows):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.LEFT
    table.style = 'Table Grid'

    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        for p in cell.paragraphs:
            p.style = doc.styles['Normal']
            for run in p.runs:
                run.bold = True
                run.font.size = Pt(10)
                run.font.name = 'Times New Roman'

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.rows[r_idx + 1].cells[c_idx]
            cell.text = str(val)
            for p in cell.paragraphs:
                p.style = doc.styles['Normal']
                for run in p.runs:
                    run.font.size = Pt(10)
                    run.font.name = 'Times New Roman'

    return table


def build_docx():
    doc = Document()

    style = doc.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = Pt(12)
    style.paragraph_format.space_after = Pt(6)
    style.paragraph_format.line_spacing = 1.15

    # Title
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("True Fruits: International Expansion")
    r.bold = True
    r.font.size = Pt(16)
    r.font.name = 'Times New Roman'

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = sub.add_run("Case Study Analysis")
    r2.font.size = Pt(12)
    r2.font.name = 'Times New Roman'
    r2.font.color.rgb = RGBColor(100, 100, 100)

    # Q1
    q1 = doc.add_heading("Q1: How Attractive Is the Smoothie Industry & How Is True Fruits Positioned?", level=2)
    for run in q1.runs:
        run.font.name = 'Times New Roman'

    h = doc.add_heading("Porter's Five Forces", level=3)
    for run in h.runs:
        run.font.name = 'Times New Roman'

    doc.add_paragraph(
        "I would say the smoothie industry is moderately attractive overall. The market is worth "
        "about $17.8 billion globally and growing around 10% a year so there is definitely growth "
        "there. But when you look at each of the forces individually things get more complicated."
    )

    doc.add_paragraph(
        "Threat of New Entrants (Moderate) - Making smoothies is not that hard but getting into "
        "stores and building refrigerated distribution costs a lot. That keeps most new competitors out.",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Supplier Power (Moderate-High) - Fruit prices change a lot with weather and seasons. "
        "True Fruits can not just use cheaper substitutes since their brand is built on no additives.",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Buyer Power (High) - Grocery stores decide what gets on the shelves and consumers can "
        "easily grab a different brand. Both of those give buyers a lot of power.",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Substitutes (High) - People can drink juice, kombucha, protein shakes, or just eat fruit "
        "instead. There are a lot of alternatives out there.",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Rivalry (High) - Innocent has Coca-Cola behind them. PepsiCo has Naked. Danone and "
        "Nestle are in the mix too. Big companies with big budgets.",
        style='List Bullet'
    )

    h2 = doc.add_heading("SWOT", level=3)
    for run in h2.runs:
        run.font.name = 'Times New Roman'

    swot_headers = ["", "Positive", "Negative"]
    swot_rows = [
        ["Internal",
         "Strengths:\n- 70%+ market share in Germany\n- Premium glass bottles, all natural\n- Bold brand identity\n- Eckes-Granini owns 67%\n- 70M euros revenue (2023)",
         "Weaknesses:\n- Only 35 employees\n- Outsources all manufacturing\n- Controversial marketing\n- Barely any intl experience\n- Glass = expensive to ship"],
        ["External",
         "Opportunities:\n- EU smoothie mkt growing 4.4%/yr\n- Health trend keeps building\n- Lots of untapped countries\n- E-commerce as new channel",
         "Threats:\n- Innocent backed by Coca-Cola\n- Fruit prices unpredictable\n- Edgy ads could backfire\n- Inflation hurts premium pricing"],
    ]
    add_simple_table(doc, swot_headers, swot_rows)
    doc.add_paragraph()

    # Q2
    q2 = doc.add_heading("Q2: How Did Bilzerian Get From 192 to 32 Countries?", level=2)
    for run in q2.runs:
        run.font.name = 'Times New Roman'

    doc.add_paragraph(
        "He used a funnel. Started with every country and kept cutting the ones that did not make "
        "sense for True Fruits."
    )
    doc.add_paragraph(
        "Round 1 (192 to about 80) - Cut countries where people cannot afford a $2.50 glass "
        "bottle smoothie. Set a minimum GDP per capita and population. Removed unstable countries too.",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Round 2 (80 to about 50) - Cut countries without cold-chain logistics since smoothies "
        "go bad without refrigeration. Also removed places without real supermarkets or with "
        "really high tariffs on imports.",
        style='List Bullet'
    )
    doc.add_paragraph(
        "Round 3 (50 to 32) - Narrowed based on proximity to Germany, whether people are into "
        "health food, and if Eckes-Granini already has contacts there.",
        style='List Bullet'
    )
    doc.add_paragraph(
        "GDP per capita and cold-chain infrastructure were probably the two biggest deal breakers. "
        "You cannot sell an expensive smoothie where people cannot pay for it and you cannot ship "
        "it somewhere it will go bad. Some other variables I think would have helped: per-capita "
        "fruit consumption, social media usage (True Fruits relies on viral stuff), and glass "
        "recycling infrastructure since their bottles are part of the brand."
    )

    # Q3
    q3 = doc.add_heading("Q3: Top Three Countries", level=2)
    for run in q3.runs:
        run.font.name = 'Times New Roman'

    doc.add_paragraph(
        "I made a scoring model to rank countries. 10 criteria, each weighted by importance. "
        "Every country scored 1-5 on each thing, multiplied by the weight, added up. Left out "
        "countries where True Fruits already sells."
    )

    score_headers = ["Country", "GDP 15%", "Pop 10%", "Mkt 15%", "Prox 10%", "Cold 10%",
                     "Retail 10%", "Health 10%", "Comp 5%", "Ease 10%", "Cult 5%", "Score"]
    score_rows = [
        ["Netherlands", "5", "3", "4", "5", "5", "5", "5", "3", "5", "5", "4.50"],
        ["UK", "4", "5", "5", "3", "5", "5", "5", "2", "4", "4", "4.30"],
        ["Denmark", "5", "2", "3", "4", "5", "5", "5", "4", "5", "5", "4.20"],
        ["Sweden", "4", "3", "3", "3", "5", "5", "5", "4", "5", "4", "4.05"],
        ["Belgium", "4", "2", "3", "5", "5", "5", "4", "3", "5", "4", "3.95"],
        ["Norway", "5", "2", "3", "3", "5", "5", "5", "4", "4", "4", "3.95"],
        ["Italy", "3", "5", "4", "4", "4", "4", "4", "3", "3", "3", "3.70"],
        ["Canada", "4", "4", "3", "1", "5", "5", "5", "3", "5", "3", "3.70"],
        ["Poland", "2", "4", "3", "4", "3", "4", "3", "4", "4", "3", "3.30"],
        ["Japan", "3", "5", "3", "1", "5", "5", "4", "3", "3", "2", "3.30"],
    ]
    add_simple_table(doc, score_headers, score_rows)
    doc.add_paragraph()

    doc.add_paragraph(
        "Formula: Score = (GDP x .15) + (Pop x .10) + (Mkt x .15) + (Prox x .10) + "
        "(Cold x .10) + (Retail x .10) + (Health x .10) + (Comp x .05) + (Ease x .10) + (Culture x .05)"
    )

    doc.add_paragraph(
        "1. Netherlands (4.50) - Right next to Germany so shipping is cheap. GDP per capita "
        "is $73K. Good grocery chains like Albert Heijn. Eckes-Granini already distributes "
        "there. EU member so no tariffs."
    )
    doc.add_paragraph(
        "2. United Kingdom (4.30) - Biggest smoothie market in Europe at 18.6% of the "
        "continent. 67 million people. Innocent is the main competitor but True Fruits has "
        "a different vibe. Brexit makes trade harder though."
    )
    doc.add_paragraph(
        "3. Denmark (4.20) - GDP per capita of $76K, one of the highest in Europe. Danes "
        "care a lot about organic and natural food. Good stepping stone into Scandinavia. "
        "EU member so no trade barriers."
    )

    # Q4
    q4 = doc.add_heading("Q4: Is Exporting the Best Mode of Entry?", level=2)
    for run in q4.runs:
        run.font.name = 'Times New Roman'

    doc.add_paragraph(
        "Exporting makes sense for now since the company is so small. But there are other "
        "options and they each have tradeoffs."
    )

    moe_headers = ["Mode", "What It Is", "Pros", "Cons", "Fit"]
    moe_rows = [
        ["Exporting", "Make in Germany, ship it", "Low cost, low risk", "Glass is heavy, spoilage", "Good for now"],
        ["Licensing", "Local co. makes your stuff", "Zero investment", "Lose quality control", "Bad fit"],
        ["Franchising", "Partner runs your system", "Fast, local knowledge", "For restaurants not CPG", "Does not apply"],
        ["Joint Venture", "Partner with local co.", "Shared risk, local help", "Share profits, conflicts", "Good for far mkts"],
        ["FDI", "Build/buy abroad", "Total control", "Way too expensive", "Too early"],
    ]
    add_simple_table(doc, moe_headers, moe_rows)
    doc.add_paragraph()

    doc.add_paragraph(
        "The criteria definitely change by mode. Exporting means proximity and cold-chain "
        "matter a ton. Joint ventures mean you care more about partner quality and IP protection. "
        "I would say start exporting to the Netherlands and Denmark, then consider a joint "
        "venture for the UK once you know the product sells there."
    )

    # Q5
    q5 = doc.add_heading("Q5: Challenges They Will Face", level=2)
    for run in q5.runs:
        run.font.name = 'Times New Roman'

    p5 = doc.add_paragraph()
    r5 = p5.add_run("Internal:")
    r5.bold = True

    doc.add_paragraph("35 people is not enough to run multiple new countries at once", style='List Bullet')
    doc.add_paragraph("Their manufacturer needs to handle way more volume", style='List Bullet')
    doc.add_paragraph("Their edgy marketing could offend people in other countries", style='List Bullet')
    doc.add_paragraph("Glass bottles are heavy, they break, and they cost a lot to ship", style='List Bullet')
    doc.add_paragraph("All the upfront costs come before any money comes in", style='List Bullet')

    p6 = doc.add_paragraph()
    r6 = p6.add_run("External:")
    r6.bold = True

    doc.add_paragraph("Getting shelf space at foreign grocery stores is super competitive", style='List Bullet')
    doc.add_paragraph("Innocent has Coca-Cola money and will fight back", style='List Bullet')
    doc.add_paragraph("Brexit means extra paperwork and possible tariffs for the UK", style='List Bullet')
    doc.add_paragraph("Different countries like different flavors and sizes", style='List Bullet')
    doc.add_paragraph("Currency changes between euros, pounds, and kroner can hurt margins", style='List Bullet')

    # Sources
    doc.add_paragraph()
    sh = doc.add_heading("Sources", level=2)
    for run in sh.runs:
        run.font.name = 'Times New Roman'

    sources = [
        'True Fruits. Wikipedia. en.wikipedia.org/wiki/True_Fruits',
        'Hulsink, W., Carvalho, L., Kunz, M., & Laker, A. "True Fruits." Rotterdam School of Management Case 813-044-1. The Case Centre / Harvard Course Pack.',
        'Duhig, A. "True Fruits Potential Viability in the Paraguay Market." Medium.',
        'Fortune Business Insights. "Smoothie Market Size and Analysis, 2026-2034."',
        'Market Data Forecast. "Europe Smoothies Market Report, 2034."',
        'IMARC Group. "Europe Smoothies Market, 2025-2033."',
        'Worldometer. "GDP per Capita, 2025." worldometers.info/gdp/gdp-per-capita/',
        'Eurostat. "Population and Population Change Statistics, 2025."',
        'Mainsights.io. "Eckes-Granini Acquires Majority Stake in True Fruits."',
        'World Bank. "B-READY 2025 Business Environment Assessment."',
        'MBA Knowledge Base. "Modes of Entry into International Business." mbaknol.com',
        'True Fruits Supplemental Spreadsheet. Harvard Course Pack.',
    ]
    for i, s in enumerate(sources, 1):
        doc.add_paragraph(f"{i}. {s}")

    doc.save(DOCX_PATH)
    print(f"Word doc saved: {DOCX_PATH}")


def build_pptx():
    prs = Presentation()
    prs.slide_width = PInches(10)
    prs.slide_height = PInches(7.5)

    def add_slide(title_text, bullets=None, table_data=None, sub=None):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

        # title box
        from pptx.util import Emu
        left = PInches(0.6)
        top = PInches(0.4)
        width = PInches(8.8)
        height = PInches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = PPt(28)
        p.font.bold = True
        p.font.color.rgb = PRGBColor(0x33, 0x33, 0x33)
        p.font.name = 'Calibri'

        content_top = PInches(1.4)
        if sub:
            sub_box = slide.shapes.add_textbox(left, PInches(1.15), width, PInches(0.4))
            stf = sub_box.text_frame
            sp = stf.paragraphs[0]
            sp.text = sub
            sp.font.size = PPt(14)
            sp.font.color.rgb = PRGBColor(0x88, 0x88, 0x88)
            sp.font.name = 'Calibri'
            content_top = PInches(1.7)

        if bullets:
            bx = slide.shapes.add_textbox(left, content_top, width, PInches(5.0))
            btf = bx.text_frame
            btf.word_wrap = True
            for i, btext in enumerate(bullets):
                if i == 0:
                    bp = btf.paragraphs[0]
                else:
                    bp = btf.add_paragraph()

                if btext.startswith("**"):
                    clean = btext.replace("**", "")
                    bp.text = clean
                    bp.font.bold = True
                    bp.font.size = PPt(16)
                    bp.font.name = 'Calibri'
                    bp.space_after = PPt(8)
                elif btext == "":
                    bp.text = ""
                    bp.space_after = PPt(6)
                else:
                    bp.text = btext
                    bp.font.size = PPt(14)
                    bp.font.name = 'Calibri'
                    bp.level = 0
                    bp.space_after = PPt(4)
                    if btext.startswith("  "):
                        bp.level = 1
                        bp.font.size = PPt(13)

        if table_data:
            rows_count = len(table_data)
            cols_count = len(table_data[0])
            tbl_width = PInches(8.5)
            col_w = tbl_width // cols_count
            tbl = slide.shapes.add_table(rows_count, cols_count,
                left, content_top, tbl_width, PInches(0.35 * rows_count)).table

            for r in range(rows_count):
                for c in range(cols_count):
                    cell = tbl.cell(r, c)
                    cell.text = str(table_data[r][c])
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = PPt(11)
                        paragraph.font.name = 'Calibri'
                        if r == 0:
                            paragraph.font.bold = True

        return slide

    # Slide 1 - Title
    add_slide("True Fruits: International Expansion",
              sub="Case Study Analysis",
              bullets=["", "", "", "[Name]", "[Class]", "[Date]"])

    # Slide 2 - Porters
    add_slide("Porter's Five Forces",
              sub="Global smoothie market: $17.8B, growing ~10%/year",
              table_data=[
                  ["Force", "Rating", "Why"],
                  ["New Entrants", "Moderate", "Cold-chain + shelf space are barriers"],
                  ["Supplier Power", "Mod-High", "Fruit prices go up and down a lot"],
                  ["Buyer Power", "High", "Retailers control shelves"],
                  ["Substitutes", "High", "Juice, kombucha, protein shakes etc"],
                  ["Rivalry", "High", "Innocent/Coca-Cola, PepsiCo, Danone"],
              ])

    # Slide 3 - SWOT
    add_slide("SWOT Analysis", table_data=[
        ["Strengths", "Weaknesses"],
        ["70%+ German market share", "35 employees"],
        ["All-natural, glass bottles", "Outsourced production"],
        ["Bold brand identity", "Controversial marketing"],
        ["Eckes-Granini backing (67%)", "Limited intl experience"],
        ["", ""],
        ["Opportunities", "Threats"],
        ["EU mkt growing 4.4%/yr", "Innocent has Coca-Cola"],
        ["Health trend keeps building", "Fruit price swings"],
        ["Untapped markets nearby", "Ads could backfire abroad"],
    ])

    # Slide 4 - Screening
    add_slide("Market Screening: 192 to 32", bullets=[
        "**How he narrowed it down:**",
        "Round 1 (192 to ~80): GDP per capita, population, political stability",
        "Round 2 (~80 to ~50): Cold-chain logistics, modern retail, tariffs",
        "Round 3 (~50 to 32): Proximity to Germany, health trends, Eckes-Granini network",
        "",
        "**Most important filters:**",
        "GDP per capita - premium product needs buyers who can afford it",
        "Cold-chain - smoothies spoil without refrigeration",
        "Proximity - keeps shipping costs down for a small company",
    ])

    # Slide 5 - Rankings
    add_slide("Country Rankings",
              sub="Weighted scoring model, 10 criteria, scored 1-5",
              table_data=[
                  ["Rank", "Country", "Score", "Main Reason"],
                  ["1", "Netherlands", "4.50", "Next to Germany, high GDP, Eckes-Granini"],
                  ["2", "UK", "4.30", "Biggest smoothie market in Europe"],
                  ["3", "Denmark", "4.20", "Rich consumers, health culture, EU"],
                  ["4", "Sweden", "4.05", "Health-conscious, good infrastructure"],
                  ["5", "Belgium", "3.95", "Right next door, EU market"],
              ])

    # Slide 6 - Why these 3
    add_slide("Why These Three?", bullets=[
        "**Netherlands**",
        "GDP/cap: $73K | 17.8M people | Borders Germany",
        "Eckes-Granini already distributes there | No tariffs",
        "",
        "**United Kingdom**",
        "Biggest smoothie market in Europe | 67M people",
        "Innocent is main competitor | Brexit adds friction",
        "",
        "**Denmark**",
        "GDP/cap: $76K | Strong organic food culture",
        "Gateway to rest of Scandinavia | EU member",
    ])

    # Slide 7 - MOE
    add_slide("Modes of Entry", table_data=[
        ["Mode", "Risk", "Cost", "Control", "Fit for True Fruits"],
        ["Exporting", "Low", "Low", "Low", "Good for now"],
        ["Licensing", "Low", "Low", "Low", "Bad - quality risk"],
        ["Franchising", "Med", "Low", "Med", "Does not apply"],
        ["Joint Venture", "Med", "Med", "Shared", "Good for far markets"],
        ["FDI", "High", "High", "Full", "Too early"],
    ])

    # Slide 8 - Challenges
    add_slide("Challenges & Recommendations", bullets=[
        "**Internal**",
        "Need to hire - 35 people is not enough",
        "Manufacturer has to handle more volume",
        "Marketing needs to be adapted for new cultures",
        "Glass bottles are heavy and fragile to ship",
        "",
        "**External**",
        "Shelf space is super competitive",
        "Innocent/Coca-Cola will push back",
        "Brexit complicates UK trade",
        "Different countries, different tastes",
        "",
        "**Plan: Go Netherlands first, then UK, then Denmark**",
    ])

    prs.save(PPTX_PATH)
    print(f"PowerPoint saved: {PPTX_PATH}")


if __name__ == "__main__":
    build_docx()
    build_pptx()
